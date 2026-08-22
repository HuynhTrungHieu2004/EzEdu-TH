import os
from typing import Iterable

import docx
import fitz  # PyMuPDF
from pptx import Presentation


def _normalize_text(parts: list[str]) -> str:
    """Joins extracted blocks and removes empty lines while preserving readability."""
    return "\n".join(part.strip() for part in parts if part and part.strip()).strip()


def extract_text_from_pdf(file_path: str) -> str:
    """Extract text from a PDF file page by page without OCR."""
    try:
        parts: list[str] = []
        with fitz.open(file_path) as pdf_document:
            for page in pdf_document:
                page_text = page.get_text("text")
                if page_text:
                    parts.append(page_text)
        return _normalize_text(parts)
    except Exception as exc:  # pragma: no cover - library-specific parsing failures
        raise ValueError("Could not read text from the PDF file.") from exc


def extract_text_pages_from_pdf_bytes(
    content: bytes,
    page_ranges: Iterable[tuple[int, int]] | None = None,
) -> list[tuple[int, str]]:
    """Extract selected one-based PDF pages while preserving page references."""
    try:
        selected = None
        if page_ranges is not None:
            selected = {
                page_number
                for start, end in page_ranges
                for page_number in range(start, end + 1)
            }
        pages: list[tuple[int, str]] = []
        with fitz.open(stream=content, filetype="pdf") as pdf_document:
            for index, page in enumerate(pdf_document):
                page_number = index + 1
                if selected is not None and page_number not in selected:
                    continue
                text = _normalize_text([page.get_text("text")])
                if text:
                    pages.append((page_number, text))
        return pages
    except Exception as exc:  # pragma: no cover - library-specific parsing failures
        raise ValueError("Could not read text from the PDF content.") from exc


def extract_text_from_docx(file_path: str) -> str:
    """Extract text from DOCX paragraphs and tables."""
    try:
        document = docx.Document(file_path)
        parts = [paragraph.text for paragraph in document.paragraphs]
        for table in document.tables:
            for row in table.rows:
                for cell in row.cells:
                    parts.append(cell.text)
        return _normalize_text(parts)
    except Exception as exc:  # pragma: no cover - library-specific parsing failures
        raise ValueError("Could not read text from the DOCX file.") from exc


def extract_text_from_pptx(file_path: str) -> str:
    """Extract text from PPTX slides, text boxes, and tables."""
    try:
        presentation = Presentation(file_path)
        parts: list[str] = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False) and shape.text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        paragraph_text = "".join(run.text for run in paragraph.runs).strip()
                        if paragraph_text:
                            parts.append(paragraph_text)
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        for cell in row.cells:
                            parts.append(cell.text)
        return _normalize_text(parts)
    except Exception as exc:  # pragma: no cover - library-specific parsing failures
        raise ValueError("Could not read text from the PPTX file.") from exc


def extract_text(file_path: str, file_type: str) -> str:
    """Extract text from PDF, DOCX, or PPTX files and ensure the result is not empty."""
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")

    if os.path.getsize(file_path) == 0:
        raise ValueError("The uploaded file is empty.")

    normalized_type = file_type.lower()
    if normalized_type == "pdf":
        text = extract_text_from_pdf(file_path)
    elif normalized_type == "docx":
        text = extract_text_from_docx(file_path)
    elif normalized_type == "pptx":
        text = extract_text_from_pptx(file_path)
    else:
        raise ValueError(f"Unsupported file type for extraction: {file_type}")

    if not text:
        raise ValueError("No readable text content could be extracted from the file.")

    return text
